"""Tests for the PowerPoint (.pptx) parser.

Builds tiny .pptx fixtures on the fly via python-pptx so we don't
ship binary blobs in the repo. Covers slide labelling, table
flattening, speaker-notes append, and graceful failure modes.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation  # type: ignore
from pptx.util import Inches  # type: ignore

from kb_parsers.pptx_parser import PptxParser


def _build_pptx(tmp_path: Path, build) -> Path:
    f = tmp_path / "fixture.pptx"
    prs = Presentation()
    build(prs)
    prs.save(str(f))
    return f


def _blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def _title_slide(prs):
    layout = prs.slide_layouts[0]  # title
    return prs.slides.add_slide(layout)


def test_pptx_parser_one_section_per_slide(tmp_path: Path) -> None:
    def build(prs):
        s1 = _title_slide(prs)
        s1.shapes.title.text = "Slide One Title"
        if s1.placeholders[1]:
            s1.placeholders[1].text = "Body of slide 1"

        s2 = _title_slide(prs)
        s2.shapes.title.text = "Slide Two Title"
        if s2.placeholders[1]:
            s2.placeholders[1].text = "Body of slide 2"

    f = _build_pptx(tmp_path, build)
    result = PptxParser().parse(f)
    assert len(result.sections) == 2
    assert [s.label for s in result.sections] == ["slide 1", "slide 2"]
    assert "Slide One Title" in result.sections[0].text
    assert "Body of slide 1" in result.sections[0].text
    assert "Slide One Title" not in result.sections[1].text


def test_pptx_parser_flattens_tables(tmp_path: Path) -> None:
    def build(prs):
        slide = _blank_slide(prs)
        # Add a 2x2 table
        rows, cols = 2, 2
        left = top = Inches(1)
        width = Inches(4)
        height = Inches(2)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.cell(0, 0).text = "h1"
        table.cell(0, 1).text = "h2"
        table.cell(1, 0).text = "v1"
        table.cell(1, 1).text = "v2"

    f = _build_pptx(tmp_path, build)
    result = PptxParser().parse(f)
    assert len(result.sections) == 1
    text = result.sections[0].text
    assert "h1\th2" in text
    assert "v1\tv2" in text
    assert "[表格]" in text


def test_pptx_parser_appends_speaker_notes(tmp_path: Path) -> None:
    def build(prs):
        slide = _title_slide(prs)
        slide.shapes.title.text = "Visible title"
        slide.notes_slide.notes_text_frame.text = "private speaker hint"

    f = _build_pptx(tmp_path, build)
    result = PptxParser().parse(f)
    assert len(result.sections) == 1
    text = result.sections[0].text
    assert "Visible title" in text
    assert "演讲者备注" in text
    assert "private speaker hint" in text


def test_pptx_parser_drops_empty_slides(tmp_path: Path) -> None:
    def build(prs):
        # Add two blank slides; neither has any text.
        _blank_slide(prs)
        _blank_slide(prs)
        s = _title_slide(prs)
        s.shapes.title.text = "Only this one has content"

    f = _build_pptx(tmp_path, build)
    result = PptxParser().parse(f)
    assert len(result.sections) == 1
    assert result.sections[0].label == "slide 3"
    assert result.metadata["slides_total"] == 3
    assert result.metadata["slides_with_text"] == 1


def test_pptx_parser_corrupt_file_returns_parse_error(tmp_path: Path) -> None:
    f = tmp_path / "broken.pptx"
    f.write_bytes(b"not actually a zip")
    result = PptxParser().parse(f)
    assert result.sections == []
    assert "parse_error" in result.metadata


def test_pptx_parser_registered_in_dispatch() -> None:
    """The dispatcher picks pptx for .pptx files. Pins the wiring
    so we don't lose it during a refactor."""
    from kb_parsers import get_parser_for

    p = get_parser_for(Path("deck.pptx"))
    assert p is not None
    assert p.parser_name == "pptx"
