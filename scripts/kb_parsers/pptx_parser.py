"""PowerPoint (.pptx) parser via python-pptx.

Walks each slide once, collects visible text from every shape that
has a text frame (title placeholder, body text, free-floating text
boxes), and flattens embedded tables to tab-separated rows. Each
slide becomes one `ParsedSection` labelled `slide N` so the
ingester / orchestrator can cite individual slides when the model
references something specific.

Speaker notes go into the slide they belong to (appended with a
"--- 演讲者备注 ---" divider) -- those are often where the real
clinical detail lives in a workshop / vendor deck.

Embedded images are ignored on purpose; the image parser handles
standalone images and a follow-up could extract embedded ones if a
clinical use case appears.
"""

from __future__ import annotations

from pathlib import Path
from typing import List

from .base import Parser, ParseResult, ParsedSection


class PptxParser(Parser):
    extensions = (".pptx",)
    parser_name = "pptx"

    def parse(self, path: Path) -> ParseResult:
        try:
            from pptx import Presentation  # type: ignore
        except ImportError:
            return ParseResult(
                sections=[],
                metadata={
                    "parser": self.parser_name,
                    "parse_error": "python-pptx not installed",
                },
            )

        try:
            prs = Presentation(str(path))
        except Exception as exc:
            return ParseResult(
                sections=[],
                metadata={
                    "parser": self.parser_name,
                    "parse_error": f"python-pptx failed to open: {exc}",
                },
            )

        sections: List[ParsedSection] = []
        for slide_index, slide in enumerate(prs.slides):
            lines: List[str] = []
            for shape in slide.shapes:
                _extract_shape_text(shape, lines)

            notes = _extract_notes(slide)
            if notes:
                lines.append("--- 演讲者备注 ---")
                lines.append(notes)

            text = "\n".join(line for line in lines if line.strip()).strip()
            if not text:
                continue

            sections.append(
                ParsedSection(
                    text=text,
                    label=f"slide {slide_index + 1}",
                    extra={"slide_index": slide_index},
                )
            )

        return ParseResult(
            sections=sections,
            metadata={
                "parser": self.parser_name,
                "slides_total": len(prs.slides),
                "slides_with_text": len(sections),
            },
        )


def _extract_shape_text(shape, lines: List[str]) -> None:
    """Append every text-bearing chunk on this shape into `lines`.

    Handles three common shape flavours:
      - shapes with a `text_frame` (text boxes, title placeholders)
      - tables (flattened to tab-separated rows)
      - group shapes (recursed)
    """
    # Group shape: recurse into children. python-pptx exposes the
    # group as iterable over .shapes; type check via attribute to
    # avoid importing the enum.
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for child in getattr(shape, "shapes", []):
            _extract_shape_text(child, lines)
        return

    if getattr(shape, "has_text_frame", False):
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text or "" for run in para.runs).strip()
            if text:
                lines.append(text)
        return

    if getattr(shape, "has_table", False):
        rows: List[str] = []
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                rows.append("\t".join(cells))
        if rows:
            lines.append("[表格]")
            lines.extend(rows)


def _extract_notes(slide) -> str:
    """Return slide-notes text, or '' when none. The notes_slide may
    be None or have an empty text_frame; both cases collapse to ''."""
    notes_slide = getattr(slide, "notes_slide", None)
    if notes_slide is None:
        return ""
    tf = getattr(notes_slide, "notes_text_frame", None)
    if tf is None:
        return ""
    return (tf.text or "").strip()
